import unittest
from logic.app_code import get_text_embedding, mistral_chat,extract_text_from_pdf,generate_pptx_from_text
from PyPDF2 import PdfReader
from pptx import Presentation
from io import BytesIO

class TestRAGUtils(unittest.TestCase):
    def test_embedding(self):
        result = get_text_embedding("Test text")
        print("\n\nTest text embeddings:")
        print(f" Text input:Test text\n Embedding:{result}")
        self.assertIsInstance(result, list)

    def test_chat(self):
        result = mistral_chat("Tell me a joke.")
        print(f"\nMistral Chat test:\n Prompt:Tell me a joke.\n Response:{result}")
        self.assertIsInstance(result, str)
        self.assertTrue(len(result.strip()) > 0)

class TestPDFTextExtraction(unittest.TestCase):
    def test_extract_text_from_real_pdf(self):
        with open("tests/test_files/sample.pdf", "rb") as f:
            text = extract_text_from_pdf(f)
        
        print("\n\nTest PDF upload and text extraction:")
        print("\n Extracted text:", text[:200])
        self.assertIsInstance(text, str)
        self.assertIn("genomics education", text.lower()) 

class TestGeneratePPT(unittest.TestCase):

    def test_generate_ppt_with_texts(self):
        slide_texts = [
            "Introduction to Python",
            "Data Types and Variables",
            "Control Structures",
            "Functions and Modules"
        ]

        ppt = generate_pptx_from_text(slide_texts)
        
        # Check type
        self.assertIsInstance(ppt, Presentation)

        # Check slide count
        self.assertEqual(len(ppt.slides), len(slide_texts))

        # Check slide content
        for i, slide in enumerate(ppt.slides):
            title = slide.shapes.title.text
            content = slide.placeholders[1].text
            self.assertEqual(title, slide_texts[i][:40])
            self.assertEqual(content, slide_texts[i])

    def test_generate_ppt_empty(self):
        slide_texts = []
        ppt = generate_pptx_from_text(slide_texts)
        self.assertEqual(len(ppt.slides), 0)


if __name__ == '__main__':
    unittest.main()
