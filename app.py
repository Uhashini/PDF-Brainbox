import streamlit as st
from PyPDF2 import PdfReader
import numpy as np
import faiss
from mistralai import Mistral
from docx import Document
from PIL import Image
import pytesseract
from pptx import Presentation
from pptx.util import Inches
import streamlit.components.v1 as components
import urllib.parse
import streamlit.components.v1 as components
import firebase_admin
from firebase_admin import credentials, firestore
import db
import os
import datetime
import json

cred_dict = json.loads(st.secrets["FIREBASE"]["service_account"])
cred = credentials.Certificate(cred_dict)

if not firebase_admin._apps:
    firebase_admin.initialize_app(cred)

db_firestore = firestore.client()

def test_firestore_connection():
    try:
        db_firestore.collection("testCollection").document("testDoc").set({
            "message": "hello from Streamlit",
            "time": datetime.datetime.now()
        })
        st.success("Firestore write successful ✅")
    except Exception as e:
        st.error(f"Firestore error: {e}")

st.set_page_config(
    page_title="PDF Brainbox",      
    page_icon="logo.png",           #favicon
    initial_sidebar_state="auto"    #sidebar behavior
)

test_firestore_connection()

query_params = st.query_params
email = query_params.get("email", [None])[0]

def fake_email(username):
    return f"{username}@pdfbrainbox.local"

# If authenticated state isn't set, default to False
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False

if "username" not in st.session_state:
    st.session_state.username = ""
    
# If email is passed from Google login, log the user in
if email and not st.session_state.authenticated:
    st.session_state.authenticated = True
    st.session_state.username = email

# Login/Signup UI
if not st.session_state.authenticated:
    st.title("Login to PDF Brainbox")

    tab1, tab2 = st.tabs([" Login", " Signup"])

    with tab1:
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        if st.button("Login"):
            # Convert username to fake email
            user_email = fake_email(username)
            if db.verify_user(user_email, password):
                st.success("Login successful!")
                st.session_state.authenticated = True
                st.session_state.username = username
                st.rerun()
            else:
                st.error("Invalid username or password")

        st.markdown("---")
        st.subheader("Or Login with Google")
        st.markdown(
            '<a href="https://uhashini.github.io/pdf-login-page/" target="_blank">'
            '<button style="padding: 0.5rem 1rem;">Login with Google</button>'
            '</a>',
            unsafe_allow_html=True
        )

    with tab2:
        new_username = st.text_input("New Username")
        new_password = st.text_input("New Password", type="password")
        if st.button("Create Account"):
            new_email = fake_email(new_username)
            if db.create_user(new_email, new_password):
                st.success("Account created! Please log in.")
            else:
                st.error("Username already exists or error creating account.")

    st.stop()

st.sidebar.write(f"👤 Logged in as: {st.session_state.username}")

if st.sidebar.button("Logout"):
    st.session_state.authenticated = False
    st.session_state.username = ""
    # Redirect to root URL
    st.markdown("""
        <meta http-equiv="refresh" content="0;url=https://pdf-brainbox.streamlit.app/" />
    """, unsafe_allow_html=True)
    st.stop()

# Load image from file
logo = "logo.png"

# Create two columns-one for logo,one for app name
col1, col2 = st.columns([2, 8])

with col1:
    st.image(logo, width=500)

with col2:
    st.markdown("<h1 style='margin: 0; padding-left: 10px;'>PDF Brainbox</h1>", unsafe_allow_html=True)

# Sidebar navigation
page = st.sidebar.selectbox("Select tools", ["Home", "Q&A", "Quiz", "Slides", "Notes","Flashcards"])

# File uploader in sidebar
uploaded_file = st.sidebar.file_uploader("Upload a PDF file", type=["pdf", "docx", "txt", "png", "jpg", "jpeg"])

def extract_text_from_file(file):
    file_type = file.name.split('.')[-1].lower()

    if file_type == 'pdf':
        pdf_reader = PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
        return text

    elif file_type == 'docx':
        doc = Document(file)
        return "\n".join([para.text for para in doc.paragraphs])

    elif file_type == 'txt':
        return file.read().decode('utf-8')

    elif file_type in ['png', 'jpg', 'jpeg']:
        pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        image = Image.open(file)
        return pytesseract.image_to_string(image)

    else:
        return ""
    
# Mistral API setup
api_key = "izPOWi2wtW6ARdOuZihOwjXja6Lzl8a0"
embed_model = "mistral-embed"
chat_model = "mistral-small-latest"
client = Mistral(api_key=api_key)

# Shared data containers
chunks = []
index = None


# Get embedding for text
def get_text_embedding(txt):
    response = client.embeddings.create(model=embed_model, inputs=[txt])
    return response.data[0].embedding


# Chat function with mistral
def mistral_chat(user_message, is_json=False):
    messages = [{"role": "user", "content": user_message}]
    if is_json:
        chat_response = client.chat.complete(
            model=chat_model,
            messages=messages,
        )
    else:
        chat_response = client.chat.complete(
            model=chat_model,
            messages=messages
        )
    return chat_response.choices[0].message.content

def log_pdf_upload(user_id, file_name):
    db_firestore.collection("userActivity").document(user_id).collection("uploads").add({
        "filename": file_name,
        "timestamp": datetime.datetime.now()
    })
def log_quiz_attempt(user_id, section, score):
    db_firestore.collection("userActivity").document(user_id).collection("quizzes").add({
        "section": section,
        "score": score,
        "timestamp": datetime.datetime.now()
    })


# HOME PAGE
if page == "Home":
    st.title("RAG Application")
    st.header("<- Upload PDF and view details")

    if uploaded_file is not None:
        # Show file info
        st.subheader("🗃️ File Information")
        st.write(f"**File Name:** {uploaded_file.name}")
        file_size_kb = len(uploaded_file.getbuffer()) / 1024
        st.write(f"**File Size:** {file_size_kb:.2f} KB")

        # Extract text from file
        text = extract_text_from_file(uploaded_file)

        if not text.strip():
            st.error("Could not extract any text from the uploaded file.")
            st.stop()

        # Show preview
        st.subheader("📄 Extracted Text (first 1000 chars)")
        st.code(text[:1000])

        # Split into chunks
        chunk_size = 512
        chunks = [text[i: i + chunk_size] for i in range(0, len(text), chunk_size)]

        # Get embeddings
        text_embeddings = np.array([get_text_embedding(chunk) for chunk in chunks])

        # Store in vector database
        d = text_embeddings.shape[1]
        index = faiss.IndexFlatL2(d)
        index.add(text_embeddings)

        # Save in session state
        st.session_state.chunks = chunks
        st.session_state.index = index
        log_pdf_upload(user_id=st.session_state.username, file_name=uploaded_file.name)
        db_firestore.collection("test").add({"msg": "Hello from Streamlit!"})

# Q&A PAGE
elif page == "Q&A":
    st.title("Document-Based Question Answering")
    st.header("Ask anything from the PDF — I’ll find the answer for you!")

    if "chunks" not in st.session_state or "index" not in st.session_state:
        st.warning("Please upload a PDF in the Home page first.")
    else:
        chunks = st.session_state.chunks
        index = st.session_state.index

        # Display chat history
        if "messages" not in st.session_state:
            st.session_state.messages = []

        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])

        # Input prompt
        if question := st.chat_input("Ask your question"):
            with st.chat_message("user"):
                st.markdown(question)
            st.session_state.messages.append({"role": "user", "content": question})

            # Get embeddings and retrieve chunks
            question_embedding = np.array([get_text_embedding(question)])
            D, I = index.search(question_embedding, k=2)
            retrieved_chunks = [chunks[i] for i in I.tolist()[0]]

            context = "\n".join(retrieved_chunks)
            prompt = f"""
                    Context information is below.
                    ---------------------
                    {context}
                    ---------------------
                    Given the context information and not prior knowledge, answer the query.
                    Query: {question}
                    Answer:
                    """
            answer = mistral_chat(prompt)

            with st.chat_message("assistant"):
                st.markdown(answer)
            st.session_state.messages.append({"role": "assistant", "content": answer})

elif page == "Quiz":
    st.title(" Quiz from PDF")

    if "chunks" not in st.session_state:
        st.warning("Please upload a PDF in the Home page first.")
    else:
        chunks = st.session_state.chunks
        full_text = " ".join(chunks)

        num_questions = st.slider("Select number of questions", 1, 5, 3)

        # Prompt for quiz generation
        quiz_prompt = f"""
You are a helpful assistant. From the following text, generate {num_questions} multiple choice questions in this JSON format:
[
  {{
    "question": "What is ...?",
    "options": ["A", "B", "C", "D"],
    "answer": "B"
  }},
  ...
]

Only return valid JSON. Do not add ``` or explanation. Here's the text:

{full_text[:3000]}
"""

        try:
            quiz_json_str = mistral_chat(quiz_prompt).strip()

            # Strip markdown if Mistral adds it
            if quiz_json_str.startswith("```json"):
                quiz_json_str = quiz_json_str.replace("```json", "").replace("```", "").strip()
            elif quiz_json_str.startswith("```"):
                quiz_json_str = quiz_json_str.replace("```", "").strip()

            import json
            quiz_data = json.loads(quiz_json_str)

        except Exception as e:
            st.error(f"!!Error generating quiz: {e}")
            st.stop()

        # Store user answers
        if "quiz_answers" not in st.session_state:
            st.session_state.quiz_answers = {}

        # Display quiz
        for i, q in enumerate(quiz_data):
            st.markdown(f"**Q{i+1}. {q['question']}**")
            user_choice = st.radio(
                label="",
                options=q["options"],
                key=f"q_{i}"
            )
            st.session_state.quiz_answers[f"q_{i}"] = {
                "selected": user_choice,
                "correct": q["answer"]
            }

        # Submit answers
        if st.button("Submit Quiz"):
            st.subheader(" Results")
            correct_count = 0
            for i in range(len(quiz_data)):
                selected = st.session_state.quiz_answers[f"q_{i}"]["selected"]
                correct = st.session_state.quiz_answers[f"q_{i}"]["correct"]
                if selected == correct:
                    st.success(f"✔️ Q{i+1}: Correct")
                    correct_count += 1
                else:
                    st.error(f"❌ Q{i+1}: Incorrect (Correct: {correct})")

            st.info(f" Score: {correct_count} / {num_questions}")

elif page=="Notes":
    st.title("Notes from PDF")
    if "chunks" not in st.session_state:
        st.warning("Upload Pdf in the sidebar first")
    else:
        chunks=st.session_state.chunks
        
        full_text="".join(chunks)
        
        prompt=f'''
        You are an expert study assistant. Create clear and concise notes from the following content.
        Content:
        {full_text[:3000]}
        Guidelines:
        - Use bullet points.
        - Summarize only the important information.
        - Do not copy exact sentences unless they are definitions or important terms.
        - Make the notes easy to review for studying.

        Output format:
        - Title (if inferred)
        - Key points (bulleted)

        Start:
        '''
        with st.spinner("Generating notes..."):
            notes = mistral_chat(prompt)

        st.markdown(notes)
        st.download_button("Download Notes", notes, file_name="study_notes.txt")

elif page=="Slides":
    st.title("Slides from PDF")
    if "chunks" not in st.session_state:
        st.warning("Upload Pdf in the sidebar first")
    else:
        chunks=st.session_state.chunks
        
        full_text="".join(chunks)
        
        prompt=f'''
        You are an expert study assistant. Create clear and concise pppt presentation slides from the following content.
        Content:
        {full_text[:3000]}
        Guidelines:
        - Use bullet points.
        - Summarize only the important information.
        - Do not copy exact sentences unless they are definitions or important terms.

        Output format:
        - Title
        - Key points (bulleted)

        Start:
        '''
        with st.spinner("Generating slides..."):
            slides = mistral_chat(prompt)

        st.markdown(slides)
        def generate_pptx_from_text(text):
            prs = Presentation()
            slide_layout = prs.slide_layouts[1]  # Title and content layout

            # Split slides by 'Title' and its bullets
            slides = text.strip().split("\n\n")
            for slide in slides:
                lines = slide.strip().split("\n")
                if not lines:
                    continue
                title = lines[0]
                content = lines[1:]

                slide_obj = prs.slides.add_slide(slide_layout)
                slide_obj.shapes.title.text = title
                body = slide_obj.placeholders[1]
                body.text = "\n".join(content)

            # Save
            pptx_path = "presentation.pptx"
            prs.save(pptx_path)
            return pptx_path

        # Your original code...
        with st.spinner("Generating slides..."):
            slides_text = mistral_chat(prompt)

        # Generate actual pptx
        pptx_file_path = generate_pptx_from_text(slides_text)

        with open(pptx_file_path, "rb") as f:
            st.download_button("Download presentation", f, file_name="presentation.pptx")

elif page == "Flashcards":
    st.title(" Flashcards from PDF")

    if "chunks" not in st.session_state:
        st.warning("Upload a PDF in the Home page first.")
    else:
        chunks = st.session_state.chunks
        full_text = " ".join(chunks)

        flashcard_prompt = f"""
You are a helpful study assistant. From the following text, create flashcards in this JSON format:

[
  {{
    "question": "What is ...?",
    "answer": "..."
  }},
  ...
]

Only return valid JSON. Do not include ``` or explanation.
Text:
{full_text[:3000]}
"""

        try:
            flashcard_json_str = mistral_chat(flashcard_prompt).strip()

            # Clean formatting if needed
            if flashcard_json_str.startswith("```json"):
                flashcard_json_str = flashcard_json_str.replace("```json", "").replace("```", "").strip()
            elif flashcard_json_str.startswith("```"):
                flashcard_json_str = flashcard_json_str.replace("```", "").strip()

            import json
            flashcards = json.loads(flashcard_json_str)

            st.session_state.flashcards = flashcards

        except Exception as e:
            st.error(f" Error generating flashcards: {e}")
            st.stop()

        # Flashcard navigation
        if "flashcard_index" not in st.session_state:
            st.session_state.flashcard_index = 0

        card = st.session_state.flashcards[st.session_state.flashcard_index]
        st.subheader(f"Flashcard {st.session_state.flashcard_index + 1} of {len(st.session_state.flashcards)}")
        st.markdown(f"**Q:** {card['question']}")
        if st.button("Show Answer"):
            st.info(f"**A:** {card['answer']}")

        col1, col2 = st.columns(2)
        with col1:
            if st.button("⬅️ Previous") and st.session_state.flashcard_index > 0:
                st.session_state.flashcard_index -= 1
                st.rerun()
        with col2:
            if st.button("➡️ Next") and st.session_state.flashcard_index < len(st.session_state.flashcards) - 1:
                st.session_state.flashcard_index += 1
                st.rerun()
