from mistralai import Mistral
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches

api_key = "izPOWi2wtW6ARdOuZihOwjXja6Lzl8a0"
embed_model = "mistral-embed"
chat_model = "mistral-small-latest"
client = Mistral(api_key=api_key)

def get_text_embedding(txt):
    response = client.embeddings.create(model=embed_model, inputs=[txt])
    return response.data[0].embedding


# Chat function with mistral
def mistral_chat(user_message, is_json=False):
    messages = [{"role": "user", "content": user_message}]
    if is_json:
        chat_response = client.chat.complete(
            model=chat_model,
            messages=messages,
        )
    else:
        chat_response = client.chat.complete(
            model=chat_model,
            messages=messages
        )
    return chat_response.choices[0].message.content

def extract_text_from_pdf(file):
    reader = PdfReader(file)
    if reader.pages:
        return reader.pages[0].extract_text()
    return ""

def generate_pptx_from_text(text):
            prs = Presentation()
            slide_layout = prs.slide_layouts[1]  # Title and content layout

            # Split slides by 'Title' and its bullets
            slides = text.strip().split("\n\n")
            for slide in slides:
                lines = slide.strip().split("\n")
                if not lines:
                    continue
                title = lines[0]
                content = lines[1:]

                slide_obj = prs.slides.add_slide(slide_layout)
                slide_obj.shapes.title.text = title
                body = slide_obj.placeholders[1]
                body.text = "\n".join(content)

            # Save
            pptx_path = "presentation.pptx"
            prs.save(pptx_path)
            return pptx_path
